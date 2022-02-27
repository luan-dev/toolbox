import os, sys, io, time

from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from pdf2image import convert_from_path

# decorator for timing
def time_it(func):
    def wrapper(*args, **kwargs):
        start_time = time.time()
        result = func(*args, **kwargs)
        print(f"{func.__name__} complete! {time.time() - start_time} seconds\n")
        return result

    return wrapper


@time_it
def convert():
    # load file
    input_file = sys.argv[1]
    pages = load_pdf(input_file)

    # create presentation
    ppt = Presentation()
    create_ppt(ppt, pages)

    # save presentation
    output_file = (
        f"{input_file.split('.', 1)[0]}.pptx"
        if len(sys.argv) < 3
        else f"{sys.argv[2].split('.', 1)[0]}.pptx"
    )
    ppt.save(output_file)

    # print file stats
    print(f"{len(pages)} pages converted to {output_file}")
    size_in_mb = round(os.stat(output_file).st_size / (1024 * 1024), 2)
    print(f"File size: {size_in_mb} MB\n")


# Load PDF
@time_it
def load_pdf(file):
    print(f"Loading {file}...")

    # this seems to be the optimal setting without losing too much quality
    pages = convert_from_path(file, dpi=150, fmt="jpeg", thread_count=4)
    return pages


# Crnate slides
@time_it
def create_ppt(ppt, pages):
    print("Creating Powerpoint...")

    # slide config
    blank_slide_layout = ppt.slide_layouts[6]
    left = top = Inches(0)
    width = ppt.slide_width
    height = ppt.slide_height

    for page in pages:
        # add slide to presentation without saving image
        with io.BytesIO() as output:
            page.save(output, format="JPEG")
            slide = ppt.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(output, left, top, width, height)


if __name__ == "__main__":
    convert()
