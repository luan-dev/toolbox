import os, io, time
import click

from pptx import Presentation
from pptx.util import Inches
from pdf2image import convert_from_path

# decorator for timing
def time_it(func):
    def wrapper(*args, **kwargs):
        start_time: float = time.time()
        result = func(*args, **kwargs)
        print(f"{func.__name__} complete! {time.time() - start_time} seconds\n")
        return result

    return wrapper

@time_it
@click.command()
@click.argument('input_file')
@click.option("-o", "--output", default=None, help="output file name (must include extension)", type=click.File('rb'))
@click.option("-l", "--legacy", default=None, help="legacy resolution to support 4:3 aspect ratio", is_flag=True)
def convert(input_file: str, output: str, legacy: bool):
    # load file
    pages: list = load_pdf(input_file)

    # create presentation
    ppt = Presentation()
    if not legacy:
        ppt.slide_width = Inches(16)
        ppt.slide_height = Inches(9)
    create_ppt(ppt, pages)

    # save presentation
    if not output:
        # replace pdf with pptx
        output = f"{input_file.split('.', 1)[0]}.pptx"
    ppt.save(output)

    # print file stats
    print(f"{len(pages)} pages converted to {output}")
    size_in_mb: int = round(os.stat(output).st_size / (1024 * 1024), 2)
    print(f"File size: {size_in_mb} MB\n")


# Load PDF
@time_it
def load_pdf(file: str):
    print(f"Loading {file}...")

    # this seems to be the optimal setting without losing too much quality
    pages = convert_from_path(file, dpi=150, fmt="jpeg", thread_count=4)
    return pages


# Crnate slides
@time_it
def create_ppt(ppt: Presentation, pages: list):
    print("Creating Powerpoint...")

    # slide config
    blank_slide_layout = ppt.slide_layouts[6]
    left = top = Inches(0)
    width = ppt.slide_width
    height = ppt.slide_height

    for page in pages:
        # add slide to presentation without saving image
        with io.BytesIO() as output:
            page.save(output, format="JPEG")
            slide = ppt.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(output, left, top, width, height)


if __name__ == "__main__":
    convert()